import os
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Directory for screenshots
SCREENSHOT_DIR = "/var/www/fleminganalytic/static/admin_manual"
BASE_URL = "https://fleminganalytic.com/stjohn/admin"

# Session ID from previous step - HARDCODED because automations are hard
SESSION_ID = "6py2nlQD6KQNfcLIT_xiz4oBdJQueBQ9bmH6DtnKqTM"

def take_screenshot(url_suffix, filename):
    full_url = f"{BASE_URL}{url_suffix}"
    output_path = os.path.join(SCREENSHOT_DIR, filename)
    print(f"Taking screenshot of {full_url} to {output_path}")
    
    # Use wkhtmltoimage with specific desktop settings
    # --width 1600 forces a wide desktop viewport
    # --crop-h 1000 captures just the "above the fold" content (Desktop Landscape)
    try:
        subprocess.run([
            "wkhtmltoimage", 
            "--quality", "90",
            "--width", "1600",
            "--crop-h", "1000",
            "--cookie", "stjohn_session", SESSION_ID,
            full_url, 
            output_path
        ], check=True, timeout=45)
    except Exception as e:
        print(f"Error capturing {full_url}: {e}")

def create_presentation():
    # Ensure directory exists
    os.makedirs(SCREENSHOT_DIR, exist_ok=True)
    
    # Capture all required screenshots
    # 1. Login
    take_screenshot("/login", "login_page.png")
    
    # 2. Dashboard
    take_screenshot("/", "dashboard.png")
    
    # 3. Services
    take_screenshot("/services", "services_list.png")
    take_screenshot("/services/new", "service_form.png")
    
    # 4. Events
    take_screenshot("/events", "events_list.png")
    take_screenshot("/events/new", "event_form.png")
    
    # 5. Ministries
    take_screenshot("/ministries", "ministries_list.png")
    
    # 6. Assets
    take_screenshot("/assets", "asset_library.png")
    
    # 7. Theme
    take_screenshot("/theme", "theme_editor.png")

    prs = Presentation()

    # Define St. John Colors
    PRIMARY_COLOR = RGBColor(30, 58, 95)    # #1e3a5f

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True

    def add_content_slide(title_text, content_text, image_filename=None):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        
        # Title at top (standard margin)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.top = Inches(0.5)
        title.left = Inches(0.5)
        title.width = Inches(9.0)
        title.height = Inches(1.0)

        # Text Box (Below Title)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(1.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(14)
        
        # Image (Below Text - Desktop Landscape View)
        if image_filename:
            img_path = os.path.join(SCREENSHOT_DIR, image_filename)
            if os.path.exists(img_path):
                # Max dimensions available (Wide and Clear)
                max_width = Inches(9.0)
                max_height = Inches(4.5)
                
                # Position
                img_left = Inches(0.5)
                img_top = Inches(2.6)
                
                # Add picture
                pic = slide.shapes.add_picture(img_path, img_left, img_top, width=max_width)
                
                # If picture is taller than allowed space, scale it
                if pic.height > max_height:
                    ratio = max_height / pic.height
                    pic.height = int(pic.height * ratio)
                    pic.width = int(pic.width * ratio)
                
                # Center horizontally
                if pic.width < max_width:
                     pic.left = int((prs.slide_width - pic.width) / 2)
                    
            else:
                # Placeholder box
                left = Inches(0.5)
                top = Inches(2.6)
                width = Inches(9.0)
                height = Inches(4.0)
                shape = slide.shapes.add_shape(1, left, top, width, height)
                shape.text_frame.text = f"Missing: {image_filename}"

    # Slide 1: Title
    add_title_slide("St. John CMS Admin Manual", "User Guide for Website Content Management")

    # Slide 2: Access & Login
    add_content_slide(
        "Accessing the Admin Panel",
        "1. Navigate to: /stjohn/admin\n"
        "2. Enter your credentials.\n\n"
        "Default Credentials:\n"
        "• Username: admin\n"
        "• Password: stjohn2024\n\n"
        "Note: You will stay logged in for 8 hours.",
        "login_page.png"
    )

    # Slide 3: Dashboard Overview
    add_content_slide(
        "Dashboard Overview",
        "The Dashboard provides a quick snapshot of your site's content.\n\n"
        "Key Features:\n"
        "• Statistics Cards: View counts for Events, Ministries, etc.\n"
        "• Recent Activity: See the latest items added.\n"
        "• Quick Actions: One-click access to add common items.\n"
        "• Navigation: Use the sidebar (left) to access all sections.",
        "dashboard.png"
    )

    # Slide 4: Worship Services
    add_content_slide(
        "Managing Worship Services",
        "Go to 'Worship Services' in the sidebar to see all services.\n\n"
        "This list shows:\n"
        "• Day & Time\n"
        "• Status (Active/Inactive)\n"
        "• Livestream Status\n\n"
        "Click 'Edit' to modify or 'Delete' to remove.",
        "services_list.png"
    )

    # Slide 5: Adding/Editing Services
    add_content_slide(
        "Service Details",
        "When adding/editing a service:\n\n"
        "• Name: E.g., 'Contemporary Service'\n"
        "• Sort Order: Lower numbers appear first.\n"
        "• Bulletin URL: Upload the weekly PDF here.\n"
        "• Livestream: Enable to show the 'Watch Live' badge.\n"
        "• Description: Add notes or details.",
        "service_form.png"
    )
    
    # Slide 6: Events List
    add_content_slide(
        "Events Overview",
        "The Events section manages all upcoming activities.\n\n"
        "• Events are displayed in cards.\n"
        "• Use 'Sign Up' links for external forms (e.g. SignUpGenius).",
        "events_list.png"
    )

    # Slide 7: Event Form
    add_content_slide(
        "Creating Events",
        "• Title: The headline of the event.\n"
        "• Image: Upload a promotional image (JPG/PNG).\n"
        "• Category: Appears as a badge (e.g. 'Youth', 'Worship').\n"
        "• Sign-Up URL: The destination for the 'Sign Up' button.\n"
        "• Details URL: Alternative for 'Learn More' link.",
        "event_form.png"
    )

    # Slide 8: Ministries
    add_content_slide(
        "Ministries",
        "Go to 'Ministries' to manage the diverse groups.\n\n"
        "• Icons: Uses Material Design icons (e.g., 'groups', 'music_note').\n"
        "• Page URL: Link to a detailed internal page or external site.",
        "ministries_list.png"
    )

    # Slide 9: Asset Library
    add_content_slide(
        "Asset Library",
        "The Asset Library allows you to manage uploaded files.\n\n"
        "• Upload Once, Use Everywhere.\n"
        "• Click on an asset to copy its URL.\n"
        "• Supports images and PDFs.",
        "asset_library.png"
    )

    # Slide 10: Liturgical Theme
    add_content_slide(
        "Liturgical Theme",
        "Change the site's color scheme.\n\n"
        "• Select a season (Advent, Lent, etc.).\n"
        "• Click 'Update Theme'.\n"
        "• Changes are applied immediately to headers and buttons.",
        "theme_editor.png"
    )

    # Save
    prs.save('/var/www/fleminganalytic/static/stjohn/St_John_CMS_Manual.pptx')
    print("Presentation saved to /var/www/fleminganalytic/static/stjohn/St_John_CMS_Manual.pptx")

if __name__ == "__main__":
    create_presentation()
