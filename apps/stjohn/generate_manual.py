from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define St. John Colors
    PRIMARY_COLOR = RGBColor(30, 58, 95)    # #1e3a5f
    SECONDARY_COLOR = RGBColor(201, 162, 39) # #c9a227

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = title_text
        subtitle.text = subtitle_text
        
        # Colorize
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True

    def add_content_slide(title_text, content_text, placeholder_text="[Insert Screenshot Here]"):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.size = Pt(32)

        # Content (Left side usually, but we'll split)
        # We'll manually position text and placeholder
        
        # Text Box (Left)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(16)
        
        # Placeholder Box for Screenshot (Right)
        left = Inches(5.2)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(3.5)
        
        shape = slide.shapes.add_shape(
            1, # msoShapeRectangle
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = PRIMARY_COLOR
        
        # Add text to shape
        shape.text_frame.text = placeholder_text
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # Slide 1: Title
    add_title_slide("St. John CMS Admin Manual", "User Guide for Website Content Management")

    # Slide 2: Access & Login
    add_content_slide(
        "Accessing the Admin Panel",
        "1. Navigate to: /stjohn/admin\n"
        "2. Enter your credentials.\n\n"
        "Default Credentials:\n"
        "• Username: admin\n"
        "• Password: stjohn2024\n\n"
        "Note: You will stay logged in for 8 hours.",
        "Insert Screenshot of Login Page"
    )

    # Slide 3: Dashboard Overview
    add_content_slide(
        "Dashboard Overview",
        "The Dashboard provides a quick snapshot of your site's content.\n\n"
        "Key Features:\n"
        "• Statistics Cards: View counts for Events, Ministries, etc.\n"
        "• Recent Activity: See the latest items added.\n"
        "• Quick Actions: One-click access to add common items.\n"
        "• Navigation: Use the sidebar (left) to access all sections.",
        "Insert Screenshot of Dashboard"
    )

    # Slide 4: Worship Services
    add_content_slide(
        "Managing Worship Services",
        "Go to 'Worship Services' in the sidebar.\n\n"
        "Adding/Editing a Service:\n"
        "• Name: E.g., 'Contemporary Service'\n"
        "• Day & Time: Displayed prominently on the site.\n"
        "• Sort Order: Controls the display order (1 = first).\n"
        "• Bulletin URL: Upload the weekly PDF here.\n"
        "• Livestream: Check this to show the 'Watch Live' badge.",
        "Insert Screenshot of Service Form"
    )

    # Slide 5: Events & Sign-Ups
    add_content_slide(
        "Events & Sign-Ups",
        "Go to 'Events & Sign Ups'.\n\n"
        "Creating an Event:\n"
        "• Title & Description: Main details.\n"
        "• Image: Upload a promotional image.\n"
        "• Category: 'Event', 'Sign-up', or 'Recurring'.\n"
        "• Sign-Up URL: Link to external forms (SignUpGenius, etc.).\n"
        "• Featured: Check this to highlight the event on the homepage.",
        "Insert Screenshot of Event Form"
    )

    # Slide 6: Ministries
    add_content_slide(
        "Ministries",
        "Go to 'Ministries' to manage the diverse groups associated with the church.\n\n"
        "• Icons: Uses Material Design icons (e.g., 'groups', 'music_note').\n"
        "• Description: Brief text explaining the ministry.\n"
        "• Page URL: Link to a detailed page or external site.",
        "Insert Screenshot of Ministries List"
    )

    # Slide 7: Asset Library
    add_content_slide(
        "Asset Library",
        "The Asset Library allows you to manage uploaded files.\n\n"
        "Features:\n"
        "• Upload Once, Use Everywhere: Upload a logo or PDF once and link it in multiple places.\n"
        "• File Types: Supports Images (JPG, PNG) and Documents (PDF).\n"
        "• Copy Links: Easily copy the URL of an uploaded file for use in emails or other content.",
        "Insert Screenshot of Asset Library"
    )

    # Slide 8: Liturgical Theme
    add_content_slide(
        "Liturgical Theme",
        "Change the entire website's color scheme to match the church season.\n\n"
        "Options:\n"
        "• Ordinary Time (Green)\n"
        "• Advent (Blue/Purple)\n"
        "• Christmas (White/Gold)\n"
        "• Lent (Purple)\n"
        "• Easter (White/Gold)\n"
        "• Pentecost (Red)\n\n"
        "This updates headers, buttons, and accents automatically.",
        "Insert Screenshot of Theme Selector"
    )

    # Save
    prs.save('/var/www/fleminganalytic/static/St_John_CMS_Manual.pptx')
    print("Presentation saved to /var/www/fleminganalytic/static/St_John_CMS_Manual.pptx")

if __name__ == "__main__":
    create_presentation()
